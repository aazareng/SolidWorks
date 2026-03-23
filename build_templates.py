"""
Builds the two starter templates:
  1. module_template.pptx  – per-module PowerPoint (Overview + Changeover sections)
  2. machine_config_template.xlsx – machine build configuration workbook
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import openpyxl
from openpyxl.styles import (
    Font, PatternFill, Alignment, Border, Side, numbers
)
from openpyxl.worksheet.datavalidation import DataValidation
from openpyxl.utils import get_column_letter
import copy
import os

# ─────────────────────────────────────────────────────────────────────────────
# HELPERS
# ─────────────────────────────────────────────────────────────────────────────

DARK_BLUE  = RGBColor(0x1F, 0x39, 0x64)   # Polypack navy
MID_BLUE   = RGBColor(0x2E, 0x74, 0xB5)
ORANGE     = RGBColor(0xFF, 0x8C, 0x00)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
RED        = RGBColor(0xC0, 0x00, 0x00)

SLIDE_W = Inches(13.33)   # Widescreen 16:9
SLIDE_H = Inches(7.5)


def add_filled_shape(slide, left, top, width, height, color, text="",
                     font_size=18, font_color=WHITE, bold=False, align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()   # no border
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.bold = bold
    return shape


def add_text_box(slide, left, top, width, height, text,
                 font_size=12, color=DARK_BLUE, bold=False,
                 align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return txBox


def set_slide_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def add_pptx_section(prs, name, slide_index):
    """Insert a named section starting at slide_index (0-based)."""
    sldIdLst = prs.slides._sldIdLst
    sectionLst = prs.presentation.find(qn('p:extLst'))
    # Build section XML the simple way via the core XML
    ext_ns = 'http://schemas.microsoft.com/office/powerpoint/2015/02/mce'
    p14_ns = 'http://schemas.microsoft.com/office/powerpoint/2010/main'

    # Use the official p:sectionLst approach
    pPr = prs.presentation
    extLst = pPr.find(qn('p:extLst'))
    if extLst is None:
        extLst = etree.SubElement(pPr, qn('p:extLst'))

    # Check for existing sectionLst ext
    section_ext = None
    for ext in extLst.findall(qn('p:ext')):
        if ext.find('{http://schemas.microsoft.com/office/powerpoint/2010/main}sectionLst') is not None:
            section_ext = ext
            break

    if section_ext is None:
        section_ext = etree.SubElement(extLst, qn('p:ext'))
        section_ext.set('uri', '{521415D9-36F7-43E2-AB2F-B90AF26B5E84}')
        sl = etree.SubElement(section_ext,
                              '{http://schemas.microsoft.com/office/powerpoint/2010/main}sectionLst')
    else:
        sl = section_ext.find('{http://schemas.microsoft.com/office/powerpoint/2010/main}sectionLst')

    import uuid
    sec_id = '{' + str(uuid.uuid4()).upper() + '}'
    sec_el = etree.SubElement(sl,
                              '{http://schemas.microsoft.com/office/powerpoint/2010/main}section')
    sec_el.set('name', name)
    sec_el.set('id', sec_id)

    # Associate slides starting at slide_index
    sldIdLst_el = prs.slides._sldIdLst
    slide_ids = list(sldIdLst_el)
    sldIdLst_sec = etree.SubElement(
        sec_el, '{http://schemas.microsoft.com/office/powerpoint/2010/main}sldIdLst')
    for sid in slide_ids[slide_index:]:
        ref = etree.SubElement(
            sldIdLst_sec,
            '{http://schemas.microsoft.com/office/powerpoint/2010/main}sldId')
        ref.set('id', sid.get('id'))

# ─────────────────────────────────────────────────────────────────────────────
# BUILD PPTX
# ─────────────────────────────────────────────────────────────────────────────

def build_pptx(path):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]   # Completely blank

    # ── SLIDE 0: Instructions / cover ────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    add_filled_shape(s, 0, 0, SLIDE_W, SLIDE_H, DARK_BLUE)
    add_text_box(s, Inches(0.5), Inches(0.4), Inches(12), Inches(1),
                 "MODULE DOCUMENTATION TEMPLATE", font_size=28, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(1.5), Inches(12), Inches(0.5),
                 "Replace this slide with your module name above. Delete this instructions slide before saving.",
                 font_size=14, color=ORANGE, align=PP_ALIGN.CENTER)

    instructions = (
        "HOW TO USE THIS TEMPLATE\n\n"
        "SECTION 1 — OVERVIEW  (slides marked  [ OV ])\n"
        "  • Slide OV-1 : Overall / hero image  →  top-left of the 2×2 grid in Section 5\n"
        "  • Slides OV-2 … OV-4 : Detail images  →  remaining cells of the 2×2 grid\n"
        "  • Add more OV slides if the module needs extra detail pages (each extra pair of 2 fills one additional grid)\n\n"
        "SECTION 2 — CHANGEOVER  (slides marked  [ CO ])\n"
        "  • One slide per changeover step image\n"
        "  • Type the procedure text for that step in the NOTES PANEL (Ctrl+Shift+H to show)\n"
        "  • The script reads notes as the numbered step text beneath each changeover image\n\n"
        "DIVIDER SLIDES  (dark slides with section titles)\n"
        "  • These are skipped by the generation script — do not delete them\n"
        "  • Notes on divider slides contain  SKIP=TRUE  — do not remove that line\n\n"
        "IMAGE GUIDELINES\n"
        "  • Paste or insert your image/screenshot to fill the large grey placeholder area\n"
        "  • Add callout arrows/shapes on top as needed — the whole slide is exported as one image\n"
        "  • Recommended: keep text minimal on slides; use the notes for procedure text\n"
    )
    add_text_box(s, Inches(0.6), Inches(2.1), Inches(12.1), Inches(5.0),
                 instructions, font_size=11, color=WHITE)
    set_slide_notes(s, "SKIP=TRUE\nInstructions slide — not exported.")

    # ── SLIDE 1: Section divider — OVERVIEW ──────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    add_filled_shape(s, 0, 0, SLIDE_W, SLIDE_H, DARK_BLUE)
    add_filled_shape(s, 0, Inches(2.8), SLIDE_W, Inches(1.9), MID_BLUE)
    add_text_box(s, Inches(0.5), Inches(3.0), Inches(12.3), Inches(1.4),
                 "SECTION 1 — OVERVIEW IMAGES", font_size=36, color=WHITE,
                 bold=True, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.6),
                 "Slides in this section → Section 5 of the manual  (2×2 grid layout)",
                 font_size=16, color=ORANGE, align=PP_ALIGN.CENTER)
    set_slide_notes(s, "SKIP=TRUE\nSection divider — Overview.")

    # ── SLIDES 2-5: Overview image slides ────────────────────────────────────
    ov_labels = [
        ("OV-1", "OVERALL VIEW", "Overall / hero image — this goes in the TOP-LEFT of the 2×2 grid.\nReplace the grey area with your image."),
        ("OV-2", "DETAIL IMAGE 2", "Detail image — TOP-RIGHT of the 2×2 grid.\nReplace the grey area with your image."),
        ("OV-3", "DETAIL IMAGE 3", "Detail image — BOTTOM-LEFT of the 2×2 grid.\nReplace the grey area with your image."),
        ("OV-4", "DETAIL IMAGE 4", "Detail image — BOTTOM-RIGHT of the 2×2 grid.\nReplace the grey area with your image."),
    ]
    for tag, title, note in ov_labels:
        s = prs.slides.add_slide(blank_layout)
        # Top bar
        add_filled_shape(s, 0, 0, SLIDE_W, Inches(0.55), DARK_BLUE)
        add_text_box(s, Inches(0.1), Inches(0.07), Inches(10), Inches(0.4),
                     title, font_size=14, color=WHITE, bold=True)
        # Corner tag — visually distinct, signals type to humans
        tag_w = Inches(1.1)
        add_filled_shape(s, SLIDE_W - tag_w, 0, tag_w, Inches(0.55),
                         ORANGE, text=f"[ {tag} ]", font_size=13,
                         font_color=DARK_BLUE, bold=True)
        # Image placeholder
        img_top  = Inches(0.62)
        img_h    = Inches(6.3)
        add_filled_shape(s, Inches(0.15), img_top, Inches(12.88) - Inches(0.3), img_h,
                         LIGHT_GRAY,
                         text="⬜  PASTE / INSERT IMAGE HERE\n\nUse Insert → Pictures, or paste a screenshot.\nAdd callout shapes/arrows on top as needed.",
                         font_size=20, font_color=RGBColor(0xA0,0xA0,0xA0))
        # Bottom strip
        add_filled_shape(s, 0, Inches(6.97), SLIDE_W, Inches(0.53),
                         LIGHT_GRAY,
                         text="Slide notes (below) are not exported — use them for step text or internal comments.",
                         font_size=10, font_color=RGBColor(0x80,0x80,0x80))
        set_slide_notes(s, f"TYPE=OVERVIEW\nTAG={tag}\n\n{note}\n\n--- Add any internal comments below this line ---")

    # ── SLIDE 6: Section divider — CHANGEOVER ────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    add_filled_shape(s, 0, 0, SLIDE_W, SLIDE_H, DARK_BLUE)
    add_filled_shape(s, 0, Inches(2.8), SLIDE_W, Inches(1.9), RED)
    add_text_box(s, Inches(0.5), Inches(3.0), Inches(12.3), Inches(1.4),
                 "SECTION 2 — CHANGEOVER IMAGES", font_size=36, color=WHITE,
                 bold=True, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.6),
                 "Slides in this section → Section 7 of the manual  |  Write step text in slide NOTES",
                 font_size=16, color=ORANGE, align=PP_ALIGN.CENTER)
    set_slide_notes(s, "SKIP=TRUE\nSection divider — Changeover.")

    # ── SLIDES 7-8: Changeover slides ────────────────────────────────────────
    for i in range(1, 3):
        s = prs.slides.add_slide(blank_layout)
        add_filled_shape(s, 0, 0, SLIDE_W, Inches(0.55), RED)
        add_text_box(s, Inches(0.1), Inches(0.07), Inches(10), Inches(0.4),
                     f"CHANGEOVER STEP {i}", font_size=14, color=WHITE, bold=True)
        tag_w = Inches(1.1)
        add_filled_shape(s, SLIDE_W - tag_w, 0, tag_w, Inches(0.55),
                         ORANGE, text=f"[ CO-{i} ]", font_size=13,
                         font_color=DARK_BLUE, bold=True)
        add_filled_shape(s, Inches(0.15), Inches(0.62), Inches(12.88) - Inches(0.3), Inches(6.3),
                         LIGHT_GRAY,
                         text="⬜  PASTE / INSERT IMAGE HERE",
                         font_size=20, font_color=RGBColor(0xA0,0xA0,0xA0))
        add_filled_shape(s, 0, Inches(6.97), SLIDE_W, Inches(0.53),
                         LIGHT_GRAY,
                         text="Write the procedure text for this step in the Notes panel below (Ctrl+Shift+H).",
                         font_size=10, font_color=RGBColor(0x80,0x80,0x80))
        set_slide_notes(
            s,
            f"TYPE=CHANGEOVER\nTAG=CO-{i}\n\n"
            f"Step {i}: [Write changeover procedure text here]\n\n"
            "Example:\n"
            "  1. Loosen the two (2) M8 bolts on the side guide bracket.\n"
            "  2. Slide the guide to the new product width dimension shown on the scale.\n"
            "  3. Re-tighten both bolts to 15 Nm."
        )

    prs.save(path)
    print(f"  Saved: {path}")


# ─────────────────────────────────────────────────────────────────────────────
# BUILD XLSX
# ─────────────────────────────────────────────────────────────────────────────

HDR_FILL   = PatternFill("solid", fgColor="1F3964")
HDR_FONT   = Font(color="FFFFFF", bold=True, name="Arial", size=11)
SUB_FILL   = PatternFill("solid", fgColor="2E74B5")
SUB_FONT   = Font(color="FFFFFF", bold=True, name="Arial", size=10)
BODY_FONT  = Font(name="Arial Narrow", size=10)
NOTE_FONT  = Font(name="Arial Narrow", size=9, italic=True, color="808080")
ORANGE_FILL = PatternFill("solid", fgColor="FFF2CC")
thin = Side(style="thin", color="AAAAAA")
BORDER = Border(left=thin, right=thin, top=thin, bottom=thin)

def hdr(ws, row, col, text, width=None, fill=HDR_FILL, font=HDR_FONT):
    c = ws.cell(row=row, column=col, value=text)
    c.fill = fill; c.font = font
    c.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
    c.border = BORDER
    if width:
        ws.column_dimensions[get_column_letter(col)].width = width
    return c

def body(ws, row, col, value="", note=False, fill=None):
    c = ws.cell(row=row, column=col, value=value)
    c.font = NOTE_FONT if note else BODY_FONT
    c.alignment = Alignment(vertical="center", wrap_text=True)
    c.border = BORDER
    if fill:
        c.fill = fill
    return c


def build_xlsx(path):
    wb = openpyxl.Workbook()

    # ── Sheet 1: MACHINE INFO ─────────────────────────────────────────────────
    ws = wb.active
    ws.title = "MACHINE INFO"
    ws.sheet_view.showGridLines = False
    ws.row_dimensions[1].height = 30

    # Title banner
    ws.merge_cells("A1:C1")
    c = ws["A1"]
    c.value = "POLYPACK — MANUAL GENERATION CONFIG"
    c.fill = HDR_FILL; c.font = Font(color="FFFFFF", bold=True, name="Arial", size=14)
    c.alignment = Alignment(horizontal="center", vertical="center")

    ws.merge_cells("A2:C2")
    c = ws["A2"]
    c.value = "Fill in all yellow cells. Cells with dropdowns will restrict your choices."
    c.font = NOTE_FONT
    c.alignment = Alignment(horizontal="center", vertical="center")
    c.fill = PatternFill("solid", fgColor="FFF2CC")

    sections = [
        ("COVER PAGE", [
            ("CUSTOMER",      "Acme Packaging Corp",  "Customer company name — appears on cover and in every page header"),
            ("MODEL",         "PS45-DL",              "Machine model number (e.g. PS45-DL, WA-LM)"),
            ("SERIAL NUMBER", "25-1234",              "Serial number (e.g. 25-1234)"),
            ("MANUAL DATE",   "March 2026",           "Month and year printed on cover (e.g. March 2026)"),
            ("YEAR (YY)",     "25",                   "Two-digit year for filename prefix (e.g. 25)"),
        ]),
        ("FILM SPECIFICATIONS", [
            ("FILM TYPE",           "Clear",          "Clear/Printed, Clear, Printed"),
            ("FILM MATERIAL",       "polyolefin",     "polyolefin, polyolefin/polyethylene, polyethylene, polyvinyl"),
            ("ROLL DIAMETER (in)",  "20",             "14, 20"),
            ("FILM WIDTH (in)",     "24",             "Enter numeric value"),
            ("FILM FORMAT",         "centerfolded",   "flat, centerfolded"),
        ]),
    ]

    row = 3
    ws.column_dimensions["A"].width = 26
    ws.column_dimensions["B"].width = 32
    ws.column_dimensions["C"].width = 48

    hdr(ws, row, 1, "FIELD", fill=SUB_FILL, font=SUB_FONT)
    hdr(ws, row, 2, "VALUE",  fill=SUB_FILL, font=SUB_FONT)
    hdr(ws, row, 3, "NOTES",  fill=SUB_FILL, font=SUB_FONT)
    row += 1

    dropdowns = {
        "FILM TYPE":          '"Clear/Printed,Clear,Printed"',
        "FILM MATERIAL":      '"polyolefin,polyolefin/polyethylene,polyethylene,polyvinyl"',
        "ROLL DIAMETER (in)": '"14,20"',
        "FILM FORMAT":        '"flat,centerfolded"',
    }

    for section_name, fields in sections:
        ws.merge_cells(f"A{row}:C{row}")
        c = ws.cell(row=row, column=1, value=section_name)
        c.fill = HDR_FILL; c.font = HDR_FONT
        c.alignment = Alignment(horizontal="left", vertical="center")
        ws.row_dimensions[row].height = 20
        row += 1

        for field, default, notes in fields:
            body(ws, row, 1, field)
            v = body(ws, row, 2, default, fill=ORANGE_FILL)
            body(ws, row, 3, notes, note=True)
            ws.row_dimensions[row].height = 18
            if field in dropdowns:
                dv = DataValidation(type="list", formula1=dropdowns[field], allow_blank=True)
                dv.sqref = f"B{row}"
                ws.add_data_validation(dv)
            row += 1

    ws.freeze_panes = "A4"

    # ── Sheet 2: MODULES ──────────────────────────────────────────────────────
    wm = wb.create_sheet("MODULES")
    wm.sheet_view.showGridLines = False

    wm.merge_cells("A1:F1")
    c = wm["A1"]
    c.value = "MODULE SELECTION & SECTION 5 DESCRIPTIONS"
    c.fill = HDR_FILL; c.font = Font(color="FFFFFF", bold=True, name="Arial", size=13)
    c.alignment = Alignment(horizontal="center", vertical="center")
    wm.row_dimensions[1].height = 28

    wm.merge_cells("A2:F2")
    c = wm["A2"]
    c.value = (
        "• INCLUDE: TRUE = section is generated, FALSE = section is skipped entirely.\n"
        "• ORDER: controls the sequence of modules in Sections 5 and 12 (lower number = earlier).\n"
        "• FOLDER NAME: must match the subfolder name inside your /modules/ directory exactly.\n"
        "• SECTION 5 DESCRIPTION: one-paragraph description inserted beneath the 2×2 image grid.\n"
        "• CHANGEOVER INCLUDE: TRUE = a changeover sub-section is generated in Section 7."
    )
    c.font = NOTE_FONT
    c.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)
    c.fill = PatternFill("solid", fgColor="FFF2CC")
    wm.row_dimensions[2].height = 72

    col_defs = [
        ("ORDER",                   6),
        ("FOLDER NAME",            28),
        ("DISPLAY NAME",           32),
        ("INCLUDE\n(Section 5+12)", 12),
        ("CHANGEOVER\nINCLUDE",    12),
        ("SECTION 5 DESCRIPTION",  60),
    ]
    for ci, (label, width) in enumerate(col_defs, 1):
        hdr(wm, 3, ci, label, width=width, fill=SUB_FILL, font=SUB_FONT)
    wm.row_dimensions[3].height = 30

    example_modules = [
        (1,  "infeed_conveyor",   "Infeed Conveyor",          "TRUE",  "TRUE",
         "The infeed conveyor transports product into the wrapping zone at a controlled rate. "
         "Belt speed is synchronized with the wrap cycle via the main drive VFD."),
        (2,  "seal_bar",          "Seal Bar Assembly",        "TRUE",  "TRUE",
         "The seal bar assembly forms the cross seal on the film. Upper and lower bars are "
         "heated to the set temperature and actuated pneumatically on each machine cycle."),
        (3,  "wrap_arm",          "Wrap Arm",                 "FALSE", "FALSE",
         "The wrap arm rotates film around the product bundle. Arm speed is controlled via "
         "the secondary VFD and is recipe-adjustable from the HMI."),
        (4,  "film_unwind",       "Film Unwind",              "TRUE",  "FALSE",
         "The film unwind supplies roll stock film to the machine. Tension is maintained by "
         "a dancer roller assembly with proximity sensor feedback."),
        (5,  "discharge_conveyor","Discharge Conveyor",       "TRUE",  "FALSE",
         "The discharge conveyor removes wrapped product from the shrink tunnel exit. "
         "Speed is adjustable from the HMI recipe screen."),
        (6,  "shrink_tunnel",     "Shrink Tunnel",            "FALSE", "FALSE",
         "The shrink tunnel applies heat to the wrapped film causing it to shrink tightly "
         "around the product. Temperature and airflow are independently adjustable."),
    ]

    inc_dv = DataValidation(type="list", formula1='"TRUE,FALSE"', allow_blank=False)
    wm.add_data_validation(inc_dv)

    for r_off, (order, folder, display, inc, co_inc, desc) in enumerate(example_modules):
        row = 4 + r_off
        body(wm, row, 1, order)
        body(wm, row, 2, folder, fill=ORANGE_FILL)
        body(wm, row, 3, display, fill=ORANGE_FILL)
        inc_cell = body(wm, row, 4, inc, fill=ORANGE_FILL)
        co_cell  = body(wm, row, 5, co_inc, fill=ORANGE_FILL)
        body(wm, row, 6, desc, fill=ORANGE_FILL)
        wm.row_dimensions[row].height = 48
        inc_dv.sqref = f"D{row}:E{row}"

    wm.freeze_panes = "A4"

    # ── Sheet 3: BOM TEMPLATE ─────────────────────────────────────────────────
    wb_bom = wb.create_sheet("BOM TEMPLATE")
    wb_bom.sheet_view.showGridLines = False

    wb_bom.merge_cells("A1:G1")
    c = wb_bom["A1"]
    c.value = "BILL OF MATERIALS — per module (duplicate this sheet per module)"
    c.fill = HDR_FILL; c.font = Font(color="FFFFFF", bold=True, name="Arial", size=12)
    c.alignment = Alignment(horizontal="center", vertical="center")
    wb_bom.row_dimensions[1].height = 26

    wb_bom.merge_cells("A2:G2")
    c = wb_bom["A2"]
    c.value = "Set FOLDER NAME below to match the module folder. This table is inserted into Section 12."
    c.font = NOTE_FONT
    c.alignment = Alignment(horizontal="left", vertical="center")
    c.fill = PatternFill("solid", fgColor="FFF2CC")

    wb_bom.cell(row=3, column=1, value="MODULE FOLDER NAME:").font = Font(bold=True, name="Arial", size=10)
    c = wb_bom.cell(row=3, column=2, value="infeed_conveyor")
    c.font = BODY_FONT; c.fill = ORANGE_FILL; c.border = BORDER

    bom_cols = [("ITEM #", 8), ("PART NUMBER", 20), ("DESCRIPTION", 40),
                ("QTY", 7), ("MATERIAL", 16), ("VENDOR", 18), ("NOTES", 30)]
    for ci, (label, width) in enumerate(bom_cols, 1):
        hdr(wb_bom, 5, ci, label, width=width, fill=SUB_FILL, font=SUB_FONT)
    wb_bom.row_dimensions[5].height = 22

    sample_bom = [
        ("1", "PP-BLT-001",  "Conveyor Belt, 18\" Wide, PU White",    "1", "Polyurethane", "Habasit",    ""),
        ("2", "PP-MTR-012",  "Drive Motor, 0.5HP 3-Phase",            "1", "Cast Iron",    "SEW",        "See VFD settings sheet"),
        ("3", "PP-BRG-004",  "Flange Bearing, 1\" Bore",              "4", "Steel",        "NSK",        "Lubricate annually"),
        ("4", "PP-BLT-002",  "V-Belt, B-Section 42\"",                "1", "Rubber",       "Gates",      ""),
    ]
    for r_off, row_data in enumerate(sample_bom):
        row = 6 + r_off
        for ci, val in enumerate(row_data, 1):
            body(wb_bom, row, ci, val, fill=ORANGE_FILL if ci > 1 else None)
        wb_bom.row_dimensions[row].height = 16

    wb_bom.freeze_panes = "A6"

    wb.save(path)
    print(f"  Saved: {path}")


# ─────────────────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    base = "/home/user/SolidWorks"
    os.makedirs(f"{base}/templates", exist_ok=True)
    os.makedirs(f"{base}/modules/infeed_conveyor", exist_ok=True)   # example structure
    os.makedirs(f"{base}/output", exist_ok=True)

    print("Building templates...")
    build_pptx(f"{base}/templates/module_template.pptx")
    build_xlsx(f"{base}/templates/machine_config_template.xlsx")
    print("Done.")
